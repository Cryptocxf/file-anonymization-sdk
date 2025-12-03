"""
PPT文件脱敏器
基于python-pptx实现PPT文件敏感信息脱敏
"""

import logging
from typing import Optional, Dict, Any

from .base_anonymizer import BaseAnonymizer
from presidio_anonymizer import AnonymizerEngine
from presidio_anonymizer.entities.engine import OperatorConfig

logger = logging.getLogger(__name__)


class PPTAnonymizer(BaseAnonymizer):
    """PPT文件脱敏器"""
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    SUPPORTED_METHODS = ['mask']
    
    def __init__(self, language: str = 'zh', verbose: bool = False, 
                 output_dir: str = "./anonymized-datas"):
        super().__init__(language, verbose, output_dir)
        
        # 初始化匿名化引擎
        self.anonymizer = AnonymizerEngine()
    
    def anonymize(self, input_path: str, output_path: Optional[str] = None, 
                 method: str = 'mask', **kwargs) -> str:
        """
        PPT文件脱敏
        
        Args:
            input_path: 输入PPT路径
            output_path: 输出路径（可选）
            method: 脱敏方法（目前只支持mask）
        
        Returns:
            str: 输出文件路径
        
        Raises:
            FileValidationError: 文件验证失败
            MethodNotSupportedError: 方法不支持
            Exception: 处理错误
        """
        from ..exceptions import FileValidationError, MethodNotSupportedError
        
        # 验证文件
        is_valid, message = self.validate_file(input_path)
        if not is_valid:
            raise FileValidationError(message)
        
        # 验证方法
        is_valid, message = self.validate_method(method)
        if not is_valid:
            raise MethodNotSupportedError(message)
        
        # 生成输出路径
        if output_path is None:
            output_path = self.get_output_path(input_path, method)
        
        logger.info(f"开始PPT脱敏: {input_path} -> {output_path}")
        logger.info(f"脱敏方法: {method}")
        
        try:
            # 获取操作符配置
            operators = self._get_mask_operators()
            
            # 执行脱敏
            total_redactions = self._anonymize_ppt(input_path, output_path, operators)
            
            logger.info(f"PPT脱敏完成，共处理 {total_redactions} 处敏感信息")
            return output_path
            
        except Exception as e:
            logger.error(f"PPT脱敏失败: {e}")
            raise
    
    def _get_mask_operators(self) -> Dict[str, OperatorConfig]:
        """获取打码操作符"""
        return {
            "PERSON": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 1)}),
            "PHONE_NUMBER": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 3)}),
            "LOCATION": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 2)}),
            "EMAIL_ADDRESS": OperatorConfig("custom", {"lambda": lambda x: self._mask_email(x)}),
            "DATE_TIME": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 4)}),
            "CREDIT_CARD": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 4)}),
            "US_BANK_NUMBER": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 4)}),
            "DEFAULT": OperatorConfig("custom", {"lambda": lambda x: self._mask_text(x, 2)}),
        }
    
    def _mask_text(self, text: str, keep_chars: int = 3) -> str:
        """打码处理"""
        if not text or len(text) <= keep_chars:
            return text
        return text[:keep_chars] + '*' * (len(text) - keep_chars)
    
    def _mask_email(self, email: str) -> str:
        """邮箱打码处理"""
        if '@' not in email:
            return self._mask_text(email, 2)
        
        local, domain = email.split('@', 1)
        masked_local = self._mask_text(local, 2)
        return f"{masked_local}@{domain}"
    
    def _anonymize_ppt(self, input_path: str, output_path: str, 
                      operators: Dict[str, OperatorConfig]) -> int:
        """执行PPT文件脱敏"""
        try:
            from pptx import Presentation
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.dml.color import RGBColor
            
            # 打开PPT文件
            prs = Presentation(input_path)
            total_redactions = 0
            
            # 遍历所有幻灯片
            for slide_num, slide in enumerate(prs.slides, 1):
                logger.debug(f"处理第 {slide_num} 张幻灯片")
                
                # 处理形状（文本框等）
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                        original_text = shape.text_frame.text
                        anonymized_text = self._process_text(original_text, operators)
                        
                        if anonymized_text != original_text:
                            total_redactions += 1
                            logger.debug(f"  发现敏感信息: '{original_text}' -> '{anonymized_text}'")
                            
                            # 保留格式
                            self._preserve_text_format(shape.text_frame, anonymized_text)
                
                # 处理表格
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                if cell.text and cell.text.strip():
                                    original_text = cell.text
                                    anonymized_text = self._process_text(original_text, operators)
                                    
                                    if anonymized_text != original_text:
                                        total_redactions += 1
                                        logger.debug(f"  表格({row_idx},{col_idx})敏感信息: '{original_text}' -> '{anonymized_text}'")
                                        
                                        # 保留格式
                                        if cell.text_frame:
                                            self._preserve_text_format(cell.text_frame, anonymized_text)
            
            # 保存脱敏后的PPT
            prs.save(output_path)
            return total_redactions
            
        except ImportError as e:
            logger.error("python-pptx库未安装，请运行: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"处理PPT失败: {e}")
            raise
    
    def _process_text(self, text: str, operators: Dict[str, OperatorConfig]) -> str:
        """处理文本"""
        try:
            # 分析敏感信息
            analyzer_results = self.analyzer.analyze(text)
            
            if analyzer_results:
                # 使用匿名化器处理
                anonymized_results = self.anonymizer.anonymize(
                    text=text,
                    analyzer_results=analyzer_results,
                    operators=operators,
                )
                return anonymized_results.text
            else:
                return text
                
        except Exception as e:
            logger.error(f"处理文本失败: {e}")
            return text
    
    def _preserve_text_format(self, text_frame, new_text: str):
        """保持文本框格式"""
        try:
            if not text_frame.paragraphs:
                return
            
            # 保存第一个段落的格式
            first_paragraph = text_frame.paragraphs[0]
            if first_paragraph.runs:
                # 保存第一个run的格式
                first_run = first_paragraph.runs[0]
                original_font = first_run.font
                
                # 安全地获取字体属性
                font_info = {
                    'name': original_font.name if hasattr(original_font, 'name') else None,
                    'size': original_font.size if hasattr(original_font, 'size') else None,
                    'bold': original_font.bold if hasattr(original_font, 'bold') else None,
                    'italic': original_font.italic if hasattr(original_font, 'italic') else None,
                }
                
                # 清除所有run的文本，但保留run结构
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                
                # 只在第一个run中添加新文本，保持原有格式
                first_run.text = new_text
                
                # 安全地恢复字体格式
                try:
                    if font_info['name']:
                        first_run.font.name = font_info['name']
                    if font_info['size']:
                        first_run.font.size = font_info['size']
                    if font_info['bold'] is not None:
                        first_run.font.bold = font_info['bold']
                    if font_info['italic'] is not None:
                        first_run.font.italic = font_info['italic']
                except Exception as e:
                    logger.warning(f"字体格式恢复失败: {e}")
            else:
                # 如果没有run，直接设置文本
                text_frame.text = new_text
                
        except Exception as e:
            logger.error(f"保持文本格式失败: {e}")
            # 如果失败，直接设置文本
            text_frame.text = new_text
    
    def extract_text(self, input_path: str) -> str:
        """
        提取PPT文本内容
        
        Args:
            input_path: PPT文件路径
        
        Returns:
            PPT文本内容
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(input_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # 提取形状文本
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                        slide_text.append(shape.text_frame.text.strip())
                
                # 提取表格文本
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text and cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if slide_text:
                    text_parts.append(f"--- 幻灯片 {slide_num} ---")
                    text_parts.extend(slide_text)
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"提取PPT文本失败: {e}")
            return ""
    
    def get_slide_count(self, input_path: str) -> int:
        """
        获取PPT幻灯片数量
        
        Args:
            input_path: PPT文件路径
        
        Returns:
            幻灯片数量
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(input_path)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"获取PPT幻灯片数量失败: {e}")
            return 0